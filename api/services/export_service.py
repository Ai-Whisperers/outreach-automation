"""
Export Service

Generates PDF and PowerPoint exports of campaign projects.
"""

import io
from datetime import datetime
from typing import Any

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import (
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
)

from ..exceptions import PDFGenerationError, PowerPointGenerationError
from ..logging_config import get_logger
from ..models import ExportResponse
from .file_operations import get_file_service

logger = get_logger("export_service")


# =============================================================================
# PDF Export Service
# =============================================================================

class PDFExportService:
    """
    Service for generating PDF exports of campaign projects.
    """

    def __init__(self):
        self.files = get_file_service()
        self.styles = getSampleStyleSheet()
        self._setup_styles()

    def _setup_styles(self):
        """Set up custom PDF styles."""
        # Title style
        self.styles.add(ParagraphStyle(
            name='CustomTitle',
            parent=self.styles['Heading1'],
            fontSize=24,
            spaceAfter=30,
            alignment=TA_CENTER,
            textColor=colors.HexColor('#1a1a1a')
        ))

        # Section header
        self.styles.add(ParagraphStyle(
            name='SectionHeader',
            parent=self.styles['Heading2'],
            fontSize=16,
            spaceBefore=20,
            spaceAfter=10,
            textColor=colors.HexColor('#2563eb')
        ))

        # Subsection
        self.styles.add(ParagraphStyle(
            name='Subsection',
            parent=self.styles['Heading3'],
            fontSize=12,
            spaceBefore=15,
            spaceAfter=8,
            textColor=colors.HexColor('#374151')
        ))

        # Body text
        self.styles.add(ParagraphStyle(
            name='CustomBody',
            parent=self.styles['Normal'],
            fontSize=10,
            leading=14,
            alignment=TA_JUSTIFY,
            spaceAfter=8
        ))

        # Bullet point
        self.styles.add(ParagraphStyle(
            name='CustomBullet',
            parent=self.styles['Normal'],
            fontSize=10,
            leftIndent=20,
            bulletIndent=10,
            spaceAfter=4
        ))

    async def export_project(
        self,
        project_id: str,
        include_research: bool = True,
        include_all_ideas: bool = False,
        top_ideas_count: int = 10
    ) -> ExportResponse:
        """
        Export project as PDF.

        Args:
            project_id: Project identifier
            include_research: Include research summary
            include_all_ideas: Include all ideas or just top ones
            top_ideas_count: Number of top ideas to include

        Returns:
            Export response with file info
        """
        logger.info(f"Generating PDF export for project: {project_id}")

        try:
            # Get project data
            metadata = self.files.get_project_metadata(project_id)

            # Create PDF buffer
            buffer = io.BytesIO()
            doc = SimpleDocTemplate(
                buffer,
                pagesize=A4,
                rightMargin=0.75*inch,
                leftMargin=0.75*inch,
                topMargin=0.75*inch,
                bottomMargin=0.75*inch
            )

            # Build content
            story = []

            # Title page
            story.extend(self._build_title_page(metadata))

            # Executive summary
            story.extend(await self._build_executive_summary(project_id))

            # Quick reference
            story.extend(await self._build_quick_reference_section(project_id))

            # Top ideas
            story.extend(await self._build_ideas_section(
                project_id,
                include_all=include_all_ideas,
                top_count=top_ideas_count
            ))

            # Research summary (optional)
            if include_research:
                story.extend(await self._build_research_section(project_id))

            # Build PDF
            doc.build(story)

            # Save to file
            pdf_content = buffer.getvalue()
            buffer.close()

            filename = f"export-{project_id}-{datetime.now().strftime('%Y%m%d')}.pdf"
            file_path = await self.files.write_file(
                project_id,
                f"exports/{filename}",
                pdf_content.decode('latin-1')  # Binary handling
            )

            # Actually save as binary
            export_path = self.files.get_project_path(project_id) / "exports" / filename
            export_path.parent.mkdir(parents=True, exist_ok=True)
            with open(export_path, 'wb') as f:
                f.write(pdf_content)

            logger.info(f"PDF export complete: {filename}")

            return ExportResponse(
                project_id=project_id,
                format="pdf",
                file_path=str(export_path),
                file_size=len(pdf_content),
                pages=None  # Would need to count pages
            )

        except Exception as e:
            logger.error(f"PDF generation failed: {e}")
            raise PDFGenerationError(str(e))

    def _build_title_page(self, metadata: dict) -> list:
        """Build title page elements."""
        elements = []

        elements.append(Spacer(1, 2*inch))

        # Project name
        elements.append(Paragraph(
            metadata.get("name", "Campaign Project"),
            self.styles['CustomTitle']
        ))

        elements.append(Spacer(1, 0.5*inch))

        # Client
        elements.append(Paragraph(
            f"Cliente: {metadata.get('client', 'N/A')}",
            self.styles['CustomBody']
        ))

        # Date
        elements.append(Paragraph(
            f"Fecha: {datetime.now().strftime('%d/%m/%Y')}",
            self.styles['CustomBody']
        ))

        # Country
        elements.append(Paragraph(
            f"País: {metadata.get('country', 'Paraguay')}",
            self.styles['CustomBody']
        ))

        elements.append(PageBreak())

        return elements

    async def _build_executive_summary(self, project_id: str) -> list:
        """Build executive summary section."""
        elements = []

        elements.append(Paragraph("Resumen Ejecutivo", self.styles['SectionHeader']))

        try:
            # Try to get brief
            brief = await self.files.read_file(project_id, "brief-original.md")
            # Extract first paragraph as summary
            lines = brief.split('\n')
            summary = ""
            for line in lines:
                if line.strip() and not line.startswith('#'):
                    summary = line.strip()
                    break

            if summary:
                elements.append(Paragraph(summary, self.styles['CustomBody']))
        except Exception:
            elements.append(Paragraph(
                "Resumen no disponible.",
                self.styles['CustomBody']
            ))

        elements.append(Spacer(1, 0.3*inch))

        return elements

    async def _build_quick_reference_section(self, project_id: str) -> list:
        """Build quick reference section."""
        elements = []

        elements.append(Paragraph("Quick Reference", self.styles['SectionHeader']))

        try:
            content = await self.files.read_file(project_id, "quick-reference.md")

            # Parse and format content
            for line in content.split('\n'):
                if line.startswith('## '):
                    elements.append(Paragraph(
                        line.replace('## ', ''),
                        self.styles['Subsection']
                    ))
                elif line.startswith('- '):
                    elements.append(Paragraph(
                        f"• {line.replace('- ', '')}",
                        self.styles['CustomBullet']
                    ))
                elif line.strip():
                    elements.append(Paragraph(line, self.styles['CustomBody']))

        except Exception:
            elements.append(Paragraph(
                "Quick reference no disponible. Ejecutar síntesis primero.",
                self.styles['CustomBody']
            ))

        elements.append(PageBreak())

        return elements

    async def _build_ideas_section(
        self,
        project_id: str,
        include_all: bool,
        top_count: int
    ) -> list:
        """Build ideas section."""
        elements = []

        elements.append(Paragraph("Ideas de Campaña", self.styles['SectionHeader']))

        idea_files = self.files.get_idea_files(project_id)

        if not idea_files:
            elements.append(Paragraph(
                "No hay ideas generadas aún.",
                self.styles['CustomBody']
            ))
            return elements

        # Get ideas and sort by score
        ideas_data = []
        for file_path in idea_files:
            try:
                content = await self.files.read_file(project_id, file_path)
                # Extract title and score
                title = "Untitled"
                score = 0

                for line in content.split('\n'):
                    if line.startswith('# '):
                        title = line.replace('# ', '').strip()
                    if 'Promedio' in line or 'Overall' in line:
                        import re
                        match = re.search(r'(\d+\.?\d*)', line)
                        if match:
                            score = float(match.group(1))

                ideas_data.append({
                    'title': title,
                    'score': score,
                    'content': content
                })
            except Exception:
                continue

        # Sort by score
        ideas_data.sort(key=lambda x: x['score'], reverse=True)

        # Limit if not including all
        if not include_all:
            ideas_data = ideas_data[:top_count]

        # Add each idea
        for i, idea in enumerate(ideas_data, 1):
            elements.append(Paragraph(
                f"{i}. {idea['title']} (Score: {idea['score']})",
                self.styles['Subsection']
            ))

            # Add brief content
            lines = idea['content'].split('\n')
            in_concept = False
            for line in lines:
                if '## Concepto' in line:
                    in_concept = True
                    continue
                if in_concept and line.startswith('##'):
                    break
                if in_concept and line.strip():
                    elements.append(Paragraph(line, self.styles['CustomBody']))

            elements.append(Spacer(1, 0.2*inch))

        elements.append(PageBreak())

        return elements

    async def _build_research_section(self, project_id: str) -> list:
        """Build research summary section."""
        elements = []

        elements.append(Paragraph("Resumen de Investigación", self.styles['SectionHeader']))

        research_files = self.files.get_research_files(project_id)

        if not research_files:
            elements.append(Paragraph(
                "No hay investigación disponible.",
                self.styles['CustomBody']
            ))
            return elements

        # Group by category
        categories = {}
        for file_path in research_files:
            parts = file_path.split('/')
            if len(parts) >= 2:
                category = parts[0]
                if category not in categories:
                    categories[category] = []
                categories[category].append(file_path)

        for category, files in sorted(categories.items()):
            elements.append(Paragraph(
                category.replace('-', ' ').title(),
                self.styles['Subsection']
            ))

            elements.append(Paragraph(
                f"{len(files)} archivo(s) de investigación",
                self.styles['CustomBullet']
            ))

        return elements


# =============================================================================
# PowerPoint Export Service
# =============================================================================

class PowerPointExportService:
    """
    Service for generating PowerPoint exports of campaign projects.
    """

    def __init__(self):
        self.files = get_file_service()

    async def export_project(
        self,
        project_id: str,
        include_research: bool = True,
        include_all_ideas: bool = False,
        top_ideas_count: int = 10
    ) -> ExportResponse:
        """
        Export project as PowerPoint.

        Args:
            project_id: Project identifier
            include_research: Include research slides
            include_all_ideas: Include all ideas
            top_ideas_count: Number of top ideas

        Returns:
            Export response with file info
        """
        logger.info(f"Generating PowerPoint export for project: {project_id}")

        try:
            # Get project data
            metadata = self.files.get_project_metadata(project_id)

            # Create presentation
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            # Title slide
            self._add_title_slide(prs, metadata)

            # Agenda
            self._add_agenda_slide(prs)

            # Quick reference
            await self._add_quick_reference_slides(prs, project_id)

            # Ideas
            await self._add_ideas_slides(
                prs, project_id,
                include_all=include_all_ideas,
                top_count=top_ideas_count
            )

            # Next steps
            self._add_next_steps_slide(prs)

            # Save presentation
            filename = f"export-{project_id}-{datetime.now().strftime('%Y%m%d')}.pptx"
            export_path = self.files.get_project_path(project_id) / "exports" / filename
            export_path.parent.mkdir(parents=True, exist_ok=True)

            prs.save(str(export_path))

            file_size = export_path.stat().st_size

            logger.info(f"PowerPoint export complete: {filename}")

            return ExportResponse(
                project_id=project_id,
                format="pptx",
                file_path=str(export_path),
                file_size=file_size,
                pages=len(prs.slides)
            )

        except Exception as e:
            logger.error(f"PowerPoint generation failed: {e}")
            raise PowerPointGenerationError(str(e))

    def _add_title_slide(self, prs: Presentation, metadata: dict):
        """Add title slide."""
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(11), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = metadata.get("name", "Campaign Project")
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER

        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(4), Inches(11), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        p = subtitle_frame.paragraphs[0]
        p.text = f"Cliente: {metadata.get('client', 'N/A')}"
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.CENTER

        # Date
        date_box = slide.shapes.add_textbox(
            Inches(1), Inches(5), Inches(11), Inches(0.5)
        )
        date_frame = date_box.text_frame
        p = date_frame.paragraphs[0]
        p.text = datetime.now().strftime('%d/%m/%Y')
        p.font.size = Pt(18)
        p.alignment = PP_ALIGN.CENTER

    def _add_agenda_slide(self, prs: Presentation):
        """Add agenda slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(12), Inches(1)
        )
        p = title_box.text_frame.paragraphs[0]
        p.text = "Agenda"
        p.font.size = Pt(36)
        p.font.bold = True

        # Items
        content_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.5), Inches(11), Inches(5)
        )
        tf = content_box.text_frame

        items = [
            "1. Quick Reference",
            "2. Top Ideas",
            "3. Scoring & Ranking",
            "4. Próximos Pasos"
        ]

        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(24)
            p.space_after = Pt(20)

    async def _add_quick_reference_slides(self, prs: Presentation, project_id: str):
        """Add quick reference slides."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(12), Inches(1)
        )
        p = title_box.text_frame.paragraphs[0]
        p.text = "Quick Reference"
        p.font.size = Pt(36)
        p.font.bold = True

        try:
            content = await self.files.read_file(project_id, "quick-reference.md")

            # Add content
            content_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.5), Inches(12), Inches(5.5)
            )
            tf = content_box.text_frame
            tf.word_wrap = True

            # Parse first few sections
            lines = content.split('\n')[:30]
            first_para = True

            for line in lines:
                if line.startswith('## '):
                    if first_para:
                        p = tf.paragraphs[0]
                        first_para = False
                    else:
                        p = tf.add_paragraph()
                    p.text = line.replace('## ', '')
                    p.font.size = Pt(18)
                    p.font.bold = True
                elif line.startswith('- '):
                    p = tf.add_paragraph()
                    p.text = f"• {line.replace('- ', '')}"
                    p.font.size = Pt(14)
                    p.level = 1

        except Exception:
            content_box = slide.shapes.add_textbox(
                Inches(1), Inches(2), Inches(11), Inches(1)
            )
            p = content_box.text_frame.paragraphs[0]
            p.text = "Quick reference no disponible"
            p.font.size = Pt(18)

    async def _add_ideas_slides(
        self,
        prs: Presentation,
        project_id: str,
        include_all: bool,
        top_count: int
    ):
        """Add idea slides."""
        idea_files = self.files.get_idea_files(project_id)

        if not idea_files:
            return

        # Get and sort ideas
        ideas_data = []
        for file_path in idea_files:
            try:
                content = await self.files.read_file(project_id, file_path)
                title = "Untitled"
                score = 0
                concept = ""

                for line in content.split('\n'):
                    if line.startswith('# '):
                        title = line.replace('# ', '').strip()
                    if 'Promedio' in line:
                        import re
                        match = re.search(r'(\d+\.?\d*)', line)
                        if match:
                            score = float(match.group(1))

                # Extract concept
                in_concept = False
                for line in content.split('\n'):
                    if '## Concepto' in line:
                        in_concept = True
                        continue
                    if in_concept and line.startswith('##'):
                        break
                    if in_concept and line.strip():
                        concept += line + " "

                ideas_data.append({
                    'title': title,
                    'score': score,
                    'concept': concept.strip()[:300]
                })
            except Exception:
                continue

        ideas_data.sort(key=lambda x: x['score'], reverse=True)

        if not include_all:
            ideas_data = ideas_data[:top_count]

        # Add slide for each idea
        for i, idea in enumerate(ideas_data, 1):
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)

            # Title with rank and score
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.5), Inches(12), Inches(1)
            )
            p = title_box.text_frame.paragraphs[0]
            p.text = f"#{i} - {idea['title']}"
            p.font.size = Pt(28)
            p.font.bold = True

            # Score badge
            score_box = slide.shapes.add_textbox(
                Inches(11), Inches(0.5), Inches(2), Inches(0.5)
            )
            p = score_box.text_frame.paragraphs[0]
            p.text = f"Score: {idea['score']}"
            p.font.size = Pt(18)
            p.font.bold = True

            # Concept
            concept_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.5), Inches(12), Inches(5)
            )
            tf = concept_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = idea['concept'] if idea['concept'] else "Concepto no disponible"
            p.font.size = Pt(16)

    def _add_next_steps_slide(self, prs: Presentation):
        """Add next steps slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(12), Inches(1)
        )
        p = title_box.text_frame.paragraphs[0]
        p.text = "Próximos Pasos"
        p.font.size = Pt(36)
        p.font.bold = True

        # Steps
        content_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.5), Inches(11), Inches(5)
        )
        tf = content_box.text_frame

        steps = [
            "1. Revisar y priorizar las ideas top",
            "2. Desarrollar concepts visuales de las 3 mejores",
            "3. Validar con el cliente",
            "4. Proceder a producción"
        ]

        for i, step in enumerate(steps):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = step
            p.font.size = Pt(24)
            p.space_after = Pt(20)


# =============================================================================
# Unified Export Service (for DI Container)
# =============================================================================

class ExportService:
    """
    Unified export service that provides both PDF and PowerPoint export capabilities.
    
    This is a wrapper class for dependency injection.
    """

    def __init__(self):
        self.pdf = PDFExportService()
        self.pptx = PowerPointExportService()

    def export_pdf(
        self,
        project_id: str,
        include_research: bool = True,
        include_all_ideas: bool = False,
        top_ideas_count: int = 10
    ) -> dict[str, Any]:
        """Export project as PDF."""
        return self.pdf.export_project(
            project_id,
            include_research,
            include_all_ideas,
            top_ideas_count
        )

    def export_powerpoint(
        self,
        project_id: str,
        include_research: bool = True,
        include_all_ideas: bool = False,
        top_ideas_count: int = 10
    ) -> dict[str, Any]:
        """Export project as PowerPoint."""
        return self.pptx.export_project(
            project_id,
            include_research,
            include_all_ideas,
            top_ideas_count
        )


# =============================================================================
# Singleton Access
# =============================================================================

_pdf_service: PDFExportService | None = None
_pptx_service: PowerPointExportService | None = None


def get_pdf_export_service() -> PDFExportService:
    """Get singleton PDF export service."""
    global _pdf_service
    if _pdf_service is None:
        _pdf_service = PDFExportService()
    return _pdf_service


def get_pptx_export_service() -> PowerPointExportService:
    """Get singleton PowerPoint export service."""
    global _pptx_service
    if _pptx_service is None:
        _pptx_service = PowerPointExportService()
    return _pptx_service
